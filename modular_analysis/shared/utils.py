"""
Utility functions for the analysis system.
"""

import os
import pandas as pd
import numpy as np
from scipy.stats import skew, kurtosis
from typing import List, Dict, Optional, Tuple
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def setup_directories(base_path: str, subdirs: List[str]) -> None:
    """Create directory structure if it doesn't exist."""
    for subdir in subdirs:
        dir_path = os.path.join(base_path, subdir)
        os.makedirs(dir_path, exist_ok=True)


def clean_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """Clean and preprocess dataframe."""
    # Make a copy to avoid modifying original
    df_clean = df.copy()
    
    # Replace zero time constants with NaN (couldn't be calculated)
    if "Time Constant (ms)" in df_clean.columns:
        df_clean["Time Constant (ms)"] = df_clean["Time Constant (ms)"].replace(0, np.nan)
    
    # Replace "None" strings with NaN for SFA columns
    sfa_cols = ["SFA10", "SFAn"]
    for col in sfa_cols:
        if col in df_clean.columns:
            df_clean[col] = df_clean[col].replace("None", np.nan)
            df_clean[col] = pd.to_numeric(df_clean[col], errors='coerce')
    
    # Remove AHP data points with 'nan' strings
    if "AHP Amplitude (mV)" in df_clean.columns:
        df_clean = df_clean.loc[~(df_clean['AHP Amplitude (mV)'] == 'nan')]
    
    return df_clean


def validate_file_exists(filepath: str) -> bool:
    """Check if file exists and is accessible."""
    try:
        return os.path.isfile(filepath) and os.access(filepath, os.R_OK)
    except Exception as e:
        logger.warning(f"Error checking file {filepath}: {e}")
        return False


def safe_divide(numerator: float, denominator: float, default: float = 0.0) -> float:
    """Safely divide two numbers, returning default if division by zero."""
    try:
        if denominator == 0:
            return default
        return numerator / denominator
    except (TypeError, ZeroDivisionError):
        return default


def format_group_label(label: str, italic: bool = False) -> str:
    """Format group label for plots, with optional italics."""
    if italic:
        return rf'$\it{{{label}}}$'
    return label


def format_factorial_label(factor1_level: str, factor2_level: str, 
                           level_italic: Dict[str, bool] = None) -> str:
    """Format factorial design label with factors on separate lines and per-level italics."""
    level_italic = level_italic or {}
    
    # Format each level with or without italics based on level name
    # Note: Don't escape underscores - matplotlib handles them fine in regular text
    if level_italic.get(factor1_level, False):
        f1 = rf'$\mathit{{{factor1_level}}}$'
    else:
        f1 = factor1_level
    
    if level_italic.get(factor2_level, False):
        f2 = rf'$\mathit{{{factor2_level}}}$'
    else:
        f2 = factor2_level
    
    # Return with newline separator
    return f"{f1}\n{f2}"


def get_measurement_categories() -> Dict[str, List[str]]:
    """Get the standard measurement categories for multiple comparison correction."""
    return {
        "Intrinsic Property": [
            "Vm (mV)",
            "Rm (MOhm)",
            "Cm (pF)",
            "Time Constant (ms)",
            "Sag",
            "Rheobase (pA)"
        ],
        "Repetitive AP Property": [
            "Max Instantaneous (Hz)",
            "Max Steady-state (Hz)",
            "ISI_CoV",
            "SFA10",
            "SFAn",
            "Initial Burst Length (ms)",
            "Maximal Burst Length (ms)"
        ],
        "Individual AP Property": [
            "AP Peak (mV)",
            "AP Threshold (mV)",
            "AP Amplitude (mV)",
            "AP Rise Time (ms)",
            "AP Half-Width (ms)",
            "APD 50 (ms)",
            "APD 90 (ms)",
            "AHP Amplitude (mV)"
        ]
    }


def categorize_measurement(measurement: str) -> str:
    """Categorize a measurement for multiple comparison correction."""
    categories = get_measurement_categories()
    
    for category, measurements in categories.items():
        if measurement in measurements:
            return category
    
    # Default category if not found
    return "Other"


# Constants for parametric/nonparametric decision
# Based on Curran, West & Finch (1996) thresholds
MIN_N_PARAMETRIC = 10  # Below this, use non-parametric (can't assess assumptions reliably)
SKEW_THRESHOLD = 2.0   # |skewness| > 2 indicates severe non-normality
KURT_THRESHOLD = 7.0   # |kurtosis| > 7 indicates severe non-normality


def should_use_parametric(data_arrays: List[np.ndarray]) -> bool:
    """
    Decide whether to use parametric or non-parametric tests based on
    sample size and distribution shape (skewness/kurtosis).
    
    This replaces Shapiro-Wilk testing, which is problematic because:
    - Overpowered for large samples (rejects normality for trivial deviations)
    - Underpowered for small samples (can't detect non-normality)
    
    Decision logic:
    1. If any group has n < MIN_N_PARAMETRIC: use non-parametric
    2. If any group has |skewness| > 2 or |kurtosis| > 7: use non-parametric
    3. Otherwise: use parametric (with Welch correction for unequal variances)
    
    Thresholds based on Curran, West & Finch (1996) guidelines.
    
    Args:
        data_arrays: List of numpy arrays, one per group
        
    Returns:
        True if parametric tests should be used, False for non-parametric
    """
    for arr in data_arrays:
        # Convert to numpy array if needed
        arr = np.asarray(arr)
        arr = arr[~np.isnan(arr)]  # Remove NaN values
        
        # Small samples: can't reliably assess assumptions
        if len(arr) < MIN_N_PARAMETRIC:
            return False
        
        # Check distribution shape (need at least 3 values for skew/kurtosis)
        if len(arr) >= 3:
            try:
                arr_skew = abs(skew(arr))
                arr_kurt = abs(kurtosis(arr))  # excess kurtosis (normal = 0)
                
                # Severe violations: use non-parametric
                if arr_skew > SKEW_THRESHOLD or arr_kurt > KURT_THRESHOLD:
                    return False
            except Exception:
                # If calculation fails, be conservative
                return False
    
    return True


def create_output_paths(base_path: str, group_names: List[str]) -> Dict[str, Dict[str, str]]:
    """Create standard output file paths for all protocols and groups."""
    paths = {}
    
    protocols = {
        'brief_current': 'afterhyperpolarization',
        'membrane_test': 'membrane_properties', 
        'gap_free': 'resting_potential',
        'current_steps': ['current_step_parameters', 'frequency_vs_current', 'attenuation']
    }
    
    for group in group_names:
        paths[group] = {}
        
        for protocol, suffix in protocols.items():
            if protocol == 'current_steps':
                # Current steps has multiple outputs
                for i, suf in enumerate(suffix):
                    key = f"{protocol}_{i+1}" if i > 0 else protocol
                    paths[group][key] = os.path.join(base_path, "Results", f"Calc_{group}_{suf}.csv")
            else:
                paths[group][protocol] = os.path.join(base_path, "Results", f"Calc_{group}_{suffix}.csv")
    
    return paths


def convert_manifest_wide_to_long(df_wide: pd.DataFrame) -> pd.DataFrame:
    """
    Convert wide format Excel manifest to long format for mixed designs.
    
    Input (wide format):
        Subject ID | condition 1: 32 | condition 2: 37 | condition 3: 42 | ...
        Scn1a_1    | file1           | file2           | file3           | ...
        WT_1       | file4           | file5           | file6           | ...
    
    Output (long format):
        Subject_ID | Group  | Condition | Filename
        Scn1a_1    | Scn1a  | 32        | file1.abf
        Scn1a_1    | Scn1a  | 37        | file2.abf
        WT_1       | WT     | 32        | file4.abf
        WT_1       | WT     | 37        | file5.abf
    
    Args:
        df_wide: Wide format DataFrame from Excel
        
    Returns:
        Long format DataFrame with Subject_ID, Group, Condition, Filename columns
    """
    # Identify subject ID column (assume first column)
    subject_col = df_wide.columns[0]
    # Exclude Subject ID and Group columns (if present) from condition columns
    condition_cols = [col for col in df_wide.columns if col != subject_col and col != 'Group']
    
    # Melt to long format
    df_long = df_wide.melt(
        id_vars=[subject_col],
        value_vars=condition_cols,
        var_name='Condition_Raw',
        value_name='Filename'
    )
    
    # Rename Subject ID column
    df_long = df_long.rename(columns={subject_col: 'Subject_ID'})
    
    # Check if Group column exists in original data
    if 'Group' in df_wide.columns:
        # Merge Group info from original wide format
        subject_to_group = df_wide.set_index(subject_col)['Group'].to_dict()
        df_long['Group'] = df_long['Subject_ID'].map(subject_to_group)
    else:
        # Extract group from Subject ID (before underscore or hyphen)
        df_long['Group'] = df_long['Subject_ID'].str.split('[_-]', n=1, expand=True)[0]
    
    # Extract condition name from column header (after colon if present)
    def extract_condition(col_name):
        """Extract condition name from column header."""
        # Ensure col_name is a string (Excel might parse it as int/float)
        col_name_str = str(col_name).strip()
        
        if ':' in col_name_str:
            # Format: "condition 1: 32" -> "32"
            condition = col_name_str.split(':', 1)[1].strip()
        else:
            # Format: "32" -> "32"
            condition = col_name_str
        
        # Normalize parentheses to underscores to match folder naming convention
        # E.g., "32(2)" -> "32_2"
        condition = condition.replace('(', '_').replace(')', '')
        
        return condition
    
    df_long['Condition'] = df_long['Condition_Raw'].apply(extract_condition)
    
    # Add .abf extension if missing
    def ensure_abf_extension(filename):
        """Ensure filename has .abf extension."""
        if pd.isna(filename):
            return filename
        filename_str = str(filename).strip()
        if not filename_str.lower().endswith('.abf'):
            return f"{filename_str}.abf"
        return filename_str
    
    df_long['Filename'] = df_long['Filename'].apply(ensure_abf_extension)
    
    # Clean up - remove intermediate columns and NaN filenames
    df_long = df_long[['Subject_ID', 'Group', 'Condition', 'Filename']]
    df_long = df_long.dropna(subset=['Filename'])
    
    logger.info(f"Converted manifest: {len(df_long)} file entries from {df_long['Subject_ID'].nunique()} subjects")
    
    return df_long


def validate_manifest(df_long: pd.DataFrame, base_path: str) -> Tuple[bool, List[str]]:
    """
    Validate manifest for completeness and consistency.
    
    Args:
        df_long: Long format manifest DataFrame
        base_path: Base directory path for finding files
        
    Returns:
        Tuple of (is_valid, error_messages)
    """
    errors = []
    
    # Check required columns
    required_cols = ['Subject_ID', 'Group', 'Condition', 'Filename']
    missing_cols = [col for col in required_cols if col not in df_long.columns]
    if missing_cols:
        errors.append(f"Missing required columns: {', '.join(missing_cols)}")
        return False, errors
    
    # Check for empty values
    for col in required_cols:
        if df_long[col].isna().any():
            n_missing = df_long[col].isna().sum()
            errors.append(f"Column '{col}' has {n_missing} missing values")
    
    # Check that each subject has same number of conditions
    conditions_per_subject = df_long.groupby('Subject_ID')['Condition'].nunique()
    expected_conditions = conditions_per_subject.mode()[0] if len(conditions_per_subject) > 0 else 0
    
    incomplete_subjects = conditions_per_subject[conditions_per_subject != expected_conditions]
    if len(incomplete_subjects) > 0:
        errors.append(f"Incomplete data: {len(incomplete_subjects)} subjects don't have all conditions. "
                     f"Expected {expected_conditions} conditions per subject.")
        for subj in incomplete_subjects.index[:5]:  # Show first 5
            n_cond = incomplete_subjects[subj]
            errors.append(f"  - {subj}: has {n_cond} conditions (expected {expected_conditions})")
    
    # Check for duplicate subject-condition combinations
    duplicates = df_long.groupby(['Subject_ID', 'Condition']).size()
    duplicates = duplicates[duplicates > 1]
    if len(duplicates) > 0:
        errors.append(f"Found {len(duplicates)} duplicate subject-condition combinations")
        for (subj, cond), count in duplicates.head(5).items():
            errors.append(f"  - {subj} at {cond}: {count} entries")
    
    # Check that each subject belongs to only one group
    groups_per_subject = df_long.groupby('Subject_ID')['Group'].nunique()
    multi_group_subjects = groups_per_subject[groups_per_subject > 1]
    if len(multi_group_subjects) > 0:
        errors.append(f"Inconsistent groups: {len(multi_group_subjects)} subjects assigned to multiple groups")
        for subj in multi_group_subjects.index[:5]:
            groups = df_long[df_long['Subject_ID'] == subj]['Group'].unique()
            errors.append(f"  - {subj}: assigned to {', '.join(groups)}")
    
    # Check that at least 2 groups and 2 conditions exist
    n_groups = df_long['Group'].nunique()
    n_conditions = df_long['Condition'].nunique()
    
    if n_groups < 2:
        errors.append(f"Need at least 2 groups for mixed design (found {n_groups})")
    if n_conditions < 2:
        errors.append(f"Need at least 2 conditions for mixed design (found {n_conditions})")
    
    is_valid = len(errors) == 0
    
    if is_valid:
        logger.info(f"Manifest validation passed: {n_groups} groups × {n_conditions} conditions = "
                   f"{n_groups * n_conditions} expected cells")
    else:
        logger.warning(f"Manifest validation failed with {len(errors)} errors")
    
    return is_valid, errors


# Category order for formatted stats table (matches standard groupings)
_STATS_TABLE_CATEGORY_ORDER = [
    "Intrinsic Property", "Repetitive AP Property", "Individual AP Property", "Other"
]


def _is_p_significant(val) -> bool:
    """True if p-value string or number is < 0.05."""
    if pd.isna(val) or val == "":
        return False
    if isinstance(val, (int, float)):
        return float(val) < 0.05
    s = str(val).strip()
    if s.startswith("<"):
        return True
    try:
        return float(s) < 0.05
    except ValueError:
        return False


def _format_p_value(p) -> str:
    """Format p-value for table display: 4 decimal places."""
    if pd.isna(p) or p is None:
        return ""
    p = float(p)
    if p < 0.0001:
        return "<0.0001"
    return f"{p:.4f}"


def _format_mean_se(mean, stderr) -> str:
    """Format mean ± SE for table display."""
    if pd.isna(mean):
        return ""
    m, s = float(mean), float(stderr) if not pd.isna(stderr) else 0
    if pd.isna(s) or s == 0:
        return f"{m:.2f}"
    # Use reasonable precision (1–2 decimal places for SE)
    se_str = f"{s:.2f}" if s >= 0.01 else f"{s:.3f}"
    return f"{m:.2f} ± {se_str}"


def create_formatted_stats_table(results_dir: str, output_path: Optional[str] = None) -> Optional[str]:
    """Create a publication-style formatted stats table from Stats_parameters.csv.

    Output: PowerPoint (.pptx) with Measurement | Group1 | Group2 | ... | Corrected Omnibus | Pairwise p-values.
    Values: mean ± SE per group. Rows grouped by category (Intrinsic, Repetitive AP, Individual AP)
    with category header rows. Plain white background, black gridlines, p<0.05 bold.

    Args:
        results_dir: Directory containing Stats_parameters.csv
        output_path: Optional path. Default: results_dir/Stats_parameters_table.pptx

    Returns:
        Path to created .pptx file, or None if failed. Requires python-pptx.
    """
    csv_path = os.path.join(results_dir, "Stats_parameters.csv")
    if not os.path.isfile(csv_path):
        return None

    try:
        df = pd.read_csv(csv_path)
    except Exception as e:
        logger.warning(f"Could not read Stats_parameters.csv: {e}")
        return None

    # Detect group columns: *_mean, *_stderr, *_n
    group_names = []
    for col in df.columns:
        if col.endswith("_mean"):
            group_names.append(col[:-5])  # strip "_mean"
    if not group_names:
        return None

    # Detect omnibus corrected column
    omnibus_corr_col = None
    for c in ["Omnibus_corrected_p", "Corrected_omnibus_p"]:
        if c in df.columns:
            omnibus_corr_col = c
            break

    # Detect pairwise corrected columns (prefer _corrected_p over _p-value)
    pairwise_cols = []
    for col in df.columns:
        if "_corrected_p" in col and col != omnibus_corr_col and "Omnibus" not in col:
            pairwise_cols.append(col)

    # Build table rows
    rows = []

    # Build output column order
    out_columns = ["Measurement"] + group_names
    if omnibus_corr_col:
        out_columns.append("Corrected Omnibus")
    pairwise_labels = [f"Corrected {c.replace('_corrected_p', '').strip()}" for c in pairwise_cols]
    out_columns.extend(pairwise_labels)

    # First row: n cells
    n_row = {"Measurement": "n cells"}
    for g in group_names:
        n_col = f"{g}_n"
        if n_col in df.columns:
            n_vals = df[n_col].dropna()
            n_val = int(n_vals.iloc[0]) if len(n_vals) > 0 else ""
            n_row[g] = str(n_val) if n_val != "" else ""
        else:
            n_row[g] = ""
    for col in out_columns[len(group_names) + 1:]:  # omnibus + pairwise
        n_row[col] = ""
    rows.append(n_row)

    # Sort measurements by category (Intrinsic, Repetitive AP, Individual AP, Other)
    df_sorted = df.copy()
    df_sorted["_sort_cat"] = df_sorted["Measurement"].apply(categorize_measurement)
    df_sorted["_sort_ord"] = df_sorted["_sort_cat"].apply(
        lambda c: _STATS_TABLE_CATEGORY_ORDER.index(c) if c in _STATS_TABLE_CATEGORY_ORDER else 99
    )
    df_sorted = df_sorted.sort_values(["_sort_ord", "Measurement"]).drop(columns=["_sort_cat", "_sort_ord"], errors="ignore")

    last_cat = None
    for _, row in df_sorted.iterrows():
        cat = categorize_measurement(row.get("Measurement", ""))
        # Insert category separator row when category changes (including before first category)
        if cat != last_cat:
            sep_row = {"Measurement": cat}
            for col in out_columns[1:]:
                sep_row[col] = ""
            rows.append(sep_row)
        last_cat = cat

        out = {"Measurement": row.get("Measurement", "")}
        for g in group_names:
            mean_col, se_col = f"{g}_mean", f"{g}_stderr"
            out[g] = _format_mean_se(row.get(mean_col), row.get(se_col))
        if omnibus_corr_col:
            out["Corrected Omnibus"] = _format_p_value(row.get(omnibus_corr_col))
        for pc, label in zip(pairwise_cols, pairwise_labels):
            out[label] = _format_p_value(row.get(pc))
        rows.append(out)

    out_df = pd.DataFrame(rows, columns=out_columns)

    pptx_path = output_path or os.path.join(results_dir, "Stats_parameters_table.pptx")

    DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    EMU_PER_INCH = 914400

    def _is_p_value_col(col: str) -> bool:
        return col == "Corrected Omnibus" or col.startswith("Corrected ")

    def _subel(parent, tag: str, **kwargs):
        el = OxmlElement(tag)
        el.attrib.update(kwargs)
        parent.append(el)
        return el

    def _center_cell(cell):
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _cell_white_black_border(cell):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for c in list(tcPr):
            tcPr.remove(c)
        sf = _subel(tcPr, "a:solidFill")
        _subel(sf, "a:srgbClr", val="FFFFFF")
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            ln = _subel(tcPr, tag, w="12700", cap="flat", cmpd="sng", algn="ctr")
            sfl = _subel(ln, "a:solidFill")
            _subel(sfl, "a:srgbClr", val="000000")
            _subel(ln, "a:prstDash", val="solid")
            _subel(ln, "a:round")
            _subel(ln, "a:headEnd", type="none", w="med", len="med")
            _subel(ln, "a:tailEnd", type="none", w="med", len="med")

    def _set_cell_margins(cell):
        try:
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.margin_left = 0
            cell.margin_right = 0
        except Exception:
            pass

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    n_rows, n_cols = len(out_df) + 1, len(out_columns)
    table_width_in = 9.5
    table_height_in = min(7.0, 0.28 * n_rows)
    table_width = Inches(table_width_in)
    table_height = Inches(table_height_in)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.35), Inches(0.5), table_width, table_height)
    table = shape.table

    # Remove table style so cell-level white fill + black borders apply (clearing caused repair)
    tbl = shape._element
    style_el = tbl.find(f".//{{{DML_NS}}}tableStyleId")
    if style_el is not None:
        style_el.getparent().remove(style_el)

    # Header row
    for c, col_name in enumerate(out_columns):
        cell = table.cell(0, c)
        cell.text = str(col_name)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10)

    category_names = set(_STATS_TABLE_CATEGORY_ORDER)
    for r, row in out_df.iterrows():
        row_idx = r + 1
        measurement = row.get("Measurement", "")

        if measurement in category_names:
            cell = table.cell(row_idx, 0)
            cell.text = measurement
            if n_cols > 1:
                cell.merge(table.cell(row_idx, n_cols - 1))
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
        else:
            for c, col in enumerate(out_columns):
                val = row.get(col, "")
                cell = table.cell(row_idx, c)
                cell.text = "" if (pd.isna(val) or val == "") else str(val)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(9)
                        r.font.bold = _is_p_value_col(col) and _is_p_significant(val)

    row_height_emu = int(table_height_in * EMU_PER_INCH / n_rows)
    col_width_emu = int(table_width_in * EMU_PER_INCH / n_cols)

    for row_idx in range(n_rows):
        for col_idx in range(n_cols):
            cell = table.cell(row_idx, col_idx)
            if not cell.is_spanned:
                _center_cell(cell)
                _set_cell_margins(cell)
                _cell_white_black_border(cell)

    # Uniform row height and column width (after content to reduce override)
    for row in table.rows:
        row.height = row_height_emu
    for col in table.columns:
        col.width = col_width_emu

    # Force row heights in XML (tr.h) - PowerPoint sometimes ignores row.height
    tbl = shape._element
    for i, tr in enumerate(tbl.findall(f".//{{{DML_NS}}}tr")):
        tr.set("h", str(row_height_emu))

    prs.save(pptx_path)
    logger.info(f"Saved formatted stats table to {pptx_path}")
    return pptx_path
